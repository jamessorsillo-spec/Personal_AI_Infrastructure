"""
DBR Account Decks Weekly Review Pipeline

Runs on Cloud Run, triggered by Cloud Scheduler every Friday 8am ET.
- Reads Google Slides from a shared Drive folder
- Reviews each deck with Claude API
- Generates a PowerPoint summary
- Uploads it to the same Drive folder
- Emails a digest with the PowerPoint attached
"""

import base64
import io
import json
import logging
import os
import smtplib
import traceback
from datetime import datetime, timezone
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

import anthropic
from flask import Flask, request
from google.cloud import secretmanager
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import google.auth

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger(__name__)

app = Flask(__name__)

# Configuration
DRIVE_FOLDER_ID = "1-9inL2uDulRh5MqFQfSU-vnCFfWH_dAc"
RECIPIENT_EMAIL = "jorsillo@tetrascience.com"
CLAUDE_MODEL = "claude-sonnet-4-20250514"
GCP_PROJECT_ID = os.environ.get("GCP_PROJECT_ID", "")


def get_secret(secret_id: str) -> str:
    """Retrieve a secret from Google Secret Manager."""
    client = secretmanager.SecretManagerServiceClient()
    name = f"projects/{GCP_PROJECT_ID}/secrets/{secret_id}/versions/latest"
    response = client.access_secret_version(request={"name": name})
    return response.payload.data.decode("UTF-8")


def get_google_services():
    """Build Google API service clients using default credentials (service account)."""
    credentials, project = google.auth.default(
        scopes=[
            "https://www.googleapis.com/auth/drive",
            "https://www.googleapis.com/auth/presentations.readonly",
            "https://www.googleapis.com/auth/gmail.send",
        ]
    )
    drive_service = build("drive", "v3", credentials=credentials)
    slides_service = build("slides", "v1", credentials=credentials)
    gmail_service = build("gmail", "v1", credentials=credentials)
    return drive_service, slides_service, gmail_service


def list_slides_in_folder(drive_service) -> list[dict]:
    """List all Google Slides presentations in the target Drive folder."""
    query = (
        f"'{DRIVE_FOLDER_ID}' in parents "
        "and mimeType='application/vnd.google-apps.presentation' "
        "and trashed=false"
    )
    results = []
    page_token = None

    while True:
        response = (
            drive_service.files()
            .list(
                q=query,
                fields="nextPageToken, files(id, name)",
                pageSize=100,
                pageToken=page_token,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
                corpora="allDrives",
            )
            .execute()
        )
        results.extend(response.get("files", []))
        page_token = response.get("nextPageToken")
        if not page_token:
            break

    logger.info("Found %d presentations in folder", len(results))
    return results


def extract_slide_text(slides_service, presentation_id: str) -> str:
    """Extract all text content from a Google Slides presentation."""
    presentation = (
        slides_service.presentations().get(presentationId=presentation_id).execute()
    )
    slides = presentation.get("slides", [])
    all_text = []

    for i, slide in enumerate(slides, 1):
        slide_texts = []
        for element in slide.get("pageElements", []):
            shape = element.get("shape", {})
            text_content = shape.get("text", {})
            for text_element in text_content.get("textElements", []):
                run = text_element.get("textRun", {})
                content = run.get("content", "").strip()
                if content:
                    slide_texts.append(content)

            # Also check tables
            table = element.get("table", {})
            for row in table.get("tableRows", []):
                for cell in row.get("tableCells", []):
                    cell_text = cell.get("text", {})
                    for text_element in cell_text.get("textElements", []):
                        run = text_element.get("textRun", {})
                        content = run.get("content", "").strip()
                        if content:
                            slide_texts.append(content)

        if slide_texts:
            all_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))

    return "\n\n".join(all_text)


def review_deck_with_claude(claude_client, deck_name: str, deck_text: str) -> dict:
    """Send deck content to Claude for review."""
    if not deck_text.strip():
        return {
            "summary": "This presentation appears to be empty or contains no extractable text.",
            "key_points": ["No text content found in the slides."],
            "feedback": "Consider adding text content to the slides.",
        }

    prompt = f"""You are reviewing a business presentation deck called "{deck_name}".

Here is the full text content extracted from the slides:

{deck_text}

Please provide the following in JSON format:
1. "summary": A 2-3 sentence executive summary of the deck's content and purpose.
2. "key_points": An array of 3-5 key points or takeaways from the deck.
3. "feedback": One piece of constructive feedback to improve the deck.

Respond ONLY with valid JSON, no markdown formatting or code blocks."""

    message = claude_client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=1024,
        messages=[{"role": "user", "content": prompt}],
    )

    response_text = message.content[0].text.strip()
    # Strip markdown code fences if present
    if response_text.startswith("```"):
        lines = response_text.split("\n")
        lines = [l for l in lines if not l.startswith("```")]
        response_text = "\n".join(lines)

    return json.loads(response_text)


def create_review_pptx(reviews: list[dict], review_date: str) -> bytes:
    """Generate a PowerPoint file with one slide per reviewed deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Cover slide ---
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DBR Account Decks - Weekly Review"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    p.alignment = PP_ALIGN.CENTER

    # Date
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"Review Date: {review_date}"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p2.alignment = PP_ALIGN.CENTER

    # Count
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = f"Decks Reviewed: {len(reviews)}"
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p3.alignment = PP_ALIGN.CENTER

    # --- One slide per deck ---
    for review in reviews:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Deck title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = review["deck_name"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)

        # Executive Summary
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p = tf.add_paragraph()
        p.text = review["summary"]
        p.font.size = Pt(12)
        p.space_after = Pt(6)

        # Key Points
        y_offset = 3.0
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(12.333), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Key Points"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        for point in review["key_points"]:
            p = tf.add_paragraph()
            p.text = f"\u2022  {point}"
            p.font.size = Pt(12)
            p.space_after = Pt(4)

        # Feedback
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Constructive Feedback"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p = tf.add_paragraph()
        p.text = review["feedback"]
        p.font.size = Pt(12)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def upload_to_drive(drive_service, file_bytes: bytes, filename: str) -> str:
    """Upload the PowerPoint file to the target Drive folder."""
    media = MediaIoBaseUpload(
        io.BytesIO(file_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=True,
    )
    file_metadata = {
        "name": filename,
        "parents": [DRIVE_FOLDER_ID],
    }
    uploaded = (
        drive_service.files()
        .create(
            body=file_metadata,
            media_body=media,
            fields="id, webViewLink",
            supportsAllDrives=True,
        )
        .execute()
    )
    logger.info("Uploaded %s (id=%s)", filename, uploaded["id"])
    return uploaded.get("webViewLink", uploaded["id"])


def send_email_with_attachment(
    gmail_service, recipient: str, subject: str, body_text: str,
    attachment_bytes: bytes, attachment_filename: str
):
    """Send an email via Gmail API with a PowerPoint attachment."""
    message = MIMEMultipart()
    message["to"] = recipient
    message["subject"] = subject

    message.attach(MIMEText(body_text, "plain"))

    part = MIMEBase("application", "octet-stream")
    part.set_payload(attachment_bytes)
    encoders.encode_base64(part)
    part.add_header("Content-Disposition", f"attachment; filename={attachment_filename}")
    message.attach(part)

    raw = base64.urlsafe_b64encode(message.as_bytes()).decode()
    gmail_service.users().messages().send(
        userId="me", body={"raw": raw}
    ).execute()
    logger.info("Email sent to %s", recipient)


def run_pipeline():
    """Main pipeline logic."""
    review_date = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    logger.info("Starting DBR review pipeline for %s", review_date)

    # Initialize clients
    api_key = get_secret("claude-api-key")
    claude_client = anthropic.Anthropic(api_key=api_key)
    drive_service, slides_service, gmail_service = get_google_services()

    # List presentations
    presentations = list_slides_in_folder(drive_service)
    if not presentations:
        logger.warning("No presentations found in folder. Exiting.")
        return {"status": "no_decks", "message": "No presentations found."}

    # Review each deck
    reviews = []
    for pres in presentations:
        pres_id = pres["id"]
        pres_name = pres["name"]
        logger.info("Processing: %s (%s)", pres_name, pres_id)

        try:
            text = extract_slide_text(slides_service, pres_id)
            logger.info("Extracted %d characters from %s", len(text), pres_name)

            result = review_deck_with_claude(claude_client, pres_name, text)
            reviews.append({
                "deck_name": pres_name,
                "summary": result["summary"],
                "key_points": result["key_points"],
                "feedback": result["feedback"],
            })
            logger.info("Review complete for %s", pres_name)
        except Exception:
            logger.error("Error processing %s: %s", pres_name, traceback.format_exc())
            reviews.append({
                "deck_name": pres_name,
                "summary": "Error: could not process this deck.",
                "key_points": ["An error occurred during processing."],
                "feedback": "N/A",
            })

    # Generate PowerPoint
    pptx_filename = f"DBR Weekly Review - {review_date}.pptx"
    pptx_bytes = create_review_pptx(reviews, review_date)
    logger.info("Generated %s (%d bytes)", pptx_filename, len(pptx_bytes))

    # Upload to Drive
    drive_link = upload_to_drive(drive_service, pptx_bytes, pptx_filename)
    logger.info("Drive link: %s", drive_link)

    # Build email body
    deck_names = "\n".join(f"  - {r['deck_name']}" for r in reviews)
    email_body = (
        f"Hi James,\n\n"
        f"Your weekly DBR Account Decks review for {review_date} is ready.\n\n"
        f"Decks reviewed ({len(reviews)}):\n{deck_names}\n\n"
        f"The PowerPoint summary is attached and has also been uploaded to the "
        f"shared Drive folder:\n{drive_link}\n\n"
        f"Best,\nDBR Review Bot"
    )

    # Send email
    send_email_with_attachment(
        gmail_service,
        RECIPIENT_EMAIL,
        f"DBR Weekly Review - {review_date}",
        email_body,
        pptx_bytes,
        pptx_filename,
    )

    return {
        "status": "success",
        "decks_reviewed": len(reviews),
        "file": pptx_filename,
        "drive_link": drive_link,
    }


@app.route("/", methods=["POST"])
def handle_trigger():
    """HTTP endpoint triggered by Cloud Scheduler."""
    try:
        result = run_pipeline()
        return json.dumps(result), 200
    except Exception:
        logger.error("Pipeline failed: %s", traceback.format_exc())
        return json.dumps({"status": "error", "message": traceback.format_exc()}), 500


@app.route("/health", methods=["GET"])
def health():
    return "ok", 200


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port)
